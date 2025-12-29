# vc-copilot-backend/app/services/file_processor.py
import io
import re
import logging
from typing import List
from pdfminer.high_level import extract_text
from pptx import Presentation

logger = logging.getLogger(__name__)

class FileProcessor:
    @staticmethod
    def extract_slides(file_content: bytes, file_type: str) -> List[str]:
        try:
            file_buffer = io.BytesIO(file_content)
            if file_type == "application/pdf":
                slides = FileProcessor._extract_from_pdf(file_buffer)
            elif file_type in ["application/vnd.openxmlformats-officedocument.presentationml.presentation", "application/vnd.ms-powerpoint"]:
                slides = FileProcessor._extract_from_pptx(file_buffer)
            else:
                raise ValueError(f"Unsupported file type: {file_type}")

            valid_slides = [slide.strip() for slide in slides if slide.strip()]
            if not valid_slides:
                raise ValueError("No readable text found in the file.")
                
            logger.info(f"Successfully extracted {len(valid_slides)} slides.")
            return valid_slides
        except Exception as e:
            logger.error(f"Error processing file: {e}")
            raise

    @staticmethod
    def _extract_from_pdf(file_buffer: io.BytesIO) -> List[str]:
        text = extract_text(file_buffer)
        return [s.strip() for s in re.split(r'\f', text) if s.strip()]

    @staticmethod
    def _extract_from_pptx(file_buffer: io.BytesIO) -> List[str]:
        prs = Presentation(file_buffer)
        return [
            "\n".join(shape.text for shape in slide.shapes if hasattr(shape, "text"))
            for slide in prs.slides
        ]